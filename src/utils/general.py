from typing import Dict
import asyncio
import io

from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import fitz
import pymupdf4llm

from aiogram.types import Message, CallbackQuery


async def get_rate_data(rate_data):
    data = f"Тариф: {rate_data['_id']}\n\
             Количество токенов: {rate_data['n_tokens']}\n\
             Количество секунд на транскрипцию: {rate_data['n_transcribed_seconds']}\n\
             Количество секунд на генерацию: {rate_data['n_generated_seconds']}\n\
             Цена: {rate_data['price']}\n\
             Тип оплаты: {rate_data['type']}"
    return data


async def is_previous_message_not_answered_yet(message: Message | CallbackQuery, tasks):
    user_id = message.from_user.id

    if user_id not in tasks:
        return False

    task = tasks[user_id]
    if isinstance(task, asyncio.Lock) and task.locked():
        text = "⏳ Пожалуйста <b>дождитесь</b> ответа на предыдущее сообщение\n"
        text += "Или вы можете завершить (/cancel) обработку запроса"
        await message.answer(text=text,
                             parse_mode="HTML")
        return True
    else:
        return False


def extract_docx_text(file_bytes):
    doc = Document(file_bytes)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])


def extract_xlsx_text(file_bytes):
    wb = load_workbook(file_bytes)
    content = ""
    for sheetname in wb.sheetnames:
        sheet = wb[sheetname]
        for row in sheet.iter_rows(values_only=True):
            content += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
    return content


def extract_pptx_text(file_bytes):
    prs = Presentation(file_bytes)
    content = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + "\n"
    return content


def extract_pdf_text(file_bytes):
    doc = fitz.open(stream=file_bytes.read(), filetype="pdf")
    to_md = pymupdf4llm.to_markdown(doc)

    return to_md
